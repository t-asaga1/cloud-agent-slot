# -*- coding: utf-8 -*-
"""プロジェクト報告用パワーポイント「スロットアプリ開発」の生成スクリプト。

使い方:
    pip install python-pptx pillow
    python3 scripts/gen_project_summary_pptx.py [出力パス]

リポジトリ内の生成済み演出画像・設定資料・アプリのスクリーンショットを使い、
これまでの作業(STEP1〜6)と今後の作業(応用)をまとめた 16:9 のスライドを出力する。
スクリーンショット(app_overview / app_playing)は /tmp/pptx_imgs/ に無ければ
該当スライドの画像を省略してテキストのみで生成する。
"""

import sys
import tempfile
from pathlib import Path

from PIL import Image, ImageDraw, ImageEnhance, ImageFilter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
IMG = ROOT / "src/assets/images"
REF = ROOT / "incoming/reference"
SCREENSHOT_DIR = Path("/tmp/pptx_imgs")

# 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# 配色(和風 × パチスロ筐体イメージ)
C_BG = RGBColor(0x14, 0x10, 0x1E)        # 濃紺黒
C_PANEL = RGBColor(0x22, 0x1C, 0x33)     # パネル
C_PANEL2 = RGBColor(0x2B, 0x24, 0x40)
C_GOLD = RGBColor(0xD9, 0xB4, 0x4A)      # 金
C_RED = RGBColor(0xC0, 0x30, 0x2E)       # 緋色
C_WHITE = RGBColor(0xF2, 0xEF, 0xE6)
C_MUTE = RGBColor(0xB9, 0xB2, 0xC8)
C_GREEN = RGBColor(0x7F, 0xC9, 0x6B)
C_AMBER = RGBColor(0xE8, 0xA8, 0x3C)

FONT = "Noto Sans CJK JP"

_tmpdir = Path(tempfile.mkdtemp(prefix="pptx_assets_"))


def _prep(path: Path, name: str, darken: float | None = None, blur: int = 0) -> Path | None:
    """webp/png を pptx 用 JPEG へ変換(必要なら暗め・ぼかし加工)。"""
    if not path.exists():
        return None
    im = Image.open(path).convert("RGB")
    if blur:
        im = im.filter(ImageFilter.GaussianBlur(blur))
    if darken is not None:
        im = ImageEnhance.Brightness(im).enhance(darken)
    out = _tmpdir / f"{name}.jpg"
    im.save(out, "JPEG", quality=88)
    return out


def _set_font(run, size, color=C_WHITE, bold=False):
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    f.name = FONT
    # 日本語(East Asian)フォントも明示指定
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", FONT)


def add_bg(slide, color=C_BG):
    shape = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)  # 1 = rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_rect(slide, x, y, w, h, color, line=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """lines = [(text, size, color, bold, indent_level), ...]"""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    first = True
    for text, size, color, bold, level in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.level = level
        p.space_after = Pt(max(2, size * 0.35))
        r = p.add_run()
        r.text = text
        _set_font(r, size, color, bold)
    return box


def add_image_fit(slide, path: Path, x, y, w, h, border=True):
    """枠 (x,y,w,h) に収まるようアスペクト維持で配置。"""
    im = Image.open(path)
    iw, ih = im.size
    scale = min(w / iw, h / ih)
    dw, dh = int(iw * scale), int(ih * scale)
    px = int(x + (w - dw) / 2)
    py = int(y + (h - dh) / 2)
    pic = slide.shapes.add_picture(str(path), px, py, dw, dh)
    if border:
        pic.line.color.rgb = C_GOLD
        pic.line.width = Pt(1.25)
    return pic


def header(slide, title, subtitle=None):
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.92), C_PANEL)
    add_rect(slide, 0, Inches(0.92), SLIDE_W, Pt(2.5), C_GOLD)
    add_rect(slide, Inches(0.35), Inches(0.18), Pt(5), Inches(0.56), C_RED)
    add_text(slide, Inches(0.55), Inches(0.08), Inches(10.5), Inches(0.8),
             [(title, 26, C_GOLD, True, 0)], anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, Inches(8.2), Inches(0.08), Inches(4.8), Inches(0.8),
                 [(subtitle, 12, C_MUTE, False, 0)], align=PP_ALIGN.RIGHT,
                 anchor=MSO_ANCHOR.MIDDLE)


def caption(slide, x, y, w, text):
    add_text(slide, x, y, w, Inches(0.3), [(text, 10.5, C_MUTE, False, 0)],
             align=PP_ALIGN.CENTER)


def bullets(items, size=14, color=C_WHITE):
    out = []
    for it in items:
        if isinstance(it, tuple):
            text, opts = it[0], it[1]
        else:
            text, opts = it, {}
        out.append((
            ("■ " if opts.get("head") else "・") + text if not opts.get("plain") else text,
            opts.get("size", size),
            opts.get("color", color),
            opts.get("bold", False),
            opts.get("level", 0),
        ))
    return out


prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def new_slide():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    return s


# ============================================================
# 素材の準備
# ============================================================
img_title_bg = _prep(IMG / "ending/ending_result_all.webp", "title_bg", darken=0.55)
img_keyart = _prep(IMG / "ending/ending_result_all.webp", "keyart")
img_battle = _prep(IMG / "battle/battle_uat_g5_double_lever.webp", "battle")
img_battle2 = _prep(IMG / "battle/battle_at_g8_stop3_shizuka_cutin.webp", "battle2")
img_yokoku = _prep(IMG / "yokoku/yokoku_shizuka_koyu1_still2_strong.webp", "yokoku")
img_kaiwa = _prep(IMG / "yokoku/yokoku_kaiwa_benkei_full.webp", "kaiwa")
faces = {
    name: _prep(REF / f"設定資料/{name}_顔.png", f"face_{name}")
    for name in ["義経", "静", "弁慶", "頼朝"]
}
img_app1 = _prep(SCREENSHOT_DIR / "app_overview.webp", "app1")
img_app2 = _prep(SCREENSHOT_DIR / "app_playing.webp", "app2")

# ============================================================
# S1: タイトル
# ============================================================
s = new_slide()
if img_title_bg:
    # 全面に一枚絵(暗め加工)
    im = Image.open(img_title_bg)
    add_image_fit(s, img_title_bg, 0, 0, SLIDE_W, SLIDE_H, border=False)
add_rect(s, 0, Inches(2.35), SLIDE_W, Inches(2.5), C_BG).fill.fore_color.rgb = C_BG
# 半透明帯の代わりに濃色帯 + 金ライン
add_rect(s, 0, Inches(2.35), SLIDE_W, Pt(2.5), C_GOLD)
add_rect(s, 0, Inches(4.85), SLIDE_W, Pt(2.5), C_GOLD)
add_text(s, 0, Inches(2.55), SLIDE_W, Inches(1.35),
         [("スロットアプリ開発", 54, C_GOLD, True, 0)],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, 0, Inches(3.85), SLIDE_W, Inches(0.9),
         [("パチスロアプリケーション「義経物語」 — これまでの作業と今後の展開", 20, C_WHITE, False, 0)],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, 0, Inches(6.9), SLIDE_W - Inches(0.4), Inches(0.4),
         [("2026-07-20 / Cloud Agent プロジェクト報告", 12, C_MUTE, False, 0)],
         align=PP_ALIGN.RIGHT)

# ============================================================
# S2: プロジェクト概要
# ============================================================
s = new_slide()
header(s, "プロジェクト概要", "義経物語")
if img_app1:
    add_image_fit(s, img_app1, Inches(0.35), Inches(1.25), Inches(6.4), Inches(5.4))
    caption(s, Inches(0.35), Inches(6.75), Inches(6.4), "アプリ全体像(筐体 UI + 液晶 + デバッグパネル)")
add_text(s, Inches(7.1), Inches(1.3), Inches(5.9), Inches(5.9), bullets([
    ("オリジナル版権「義経物語」のパチスロを 1 からアプリ化", {"head": True, "bold": True, "size": 16}),
    ("Excel 仕様書を SPEC.md へ全転記し、確定事項 43 件を積み上げて仕様を確定", {}),
    ("リール制御(20 コマ × 3 リール・5 ライン・4 コマスベリ)/ 出玉抽せん / AT・上位 AT / 演出システムまで実機同等のゲームフローを実装", {}),
    ("Web(ブラウザ)と Windows デスクトップ(Tauri 2.x / exe 配布)の 2 形態", {}),
    ("品質の担保", {"head": True, "bold": True, "size": 16}),
    ("自動テスト 434 件(停止制御の全押下位置網羅・抽せんテーブル検算など)", {}),
    ("100 万 G シミュレーションで出玉検証(機械割 約 129〜131% でユーザー承認済み)", {}),
    ("演出素材は AI 生成画像 + ユーザー入稿の実素材(背景動画・BGM・SE)", {}),
], size=13.5))

# ============================================================
# S3: 開発環境
# ============================================================
s = new_slide()
header(s, "開発環境", "技術スタックと AI 素材生成")
col_w = Inches(6.15)
add_rect(s, Inches(0.35), Inches(1.25), col_w, Inches(5.7), C_PANEL)
add_text(s, Inches(0.6), Inches(1.4), col_w - Inches(0.5), Inches(5.4), bullets([
    ("アプリ実装", {"head": True, "bold": True, "size": 16, "color": C_GOLD}),
    ("TypeScript + React + Vite(Web)/ Vitest(テスト 434 件)/ oxlint", {}),
    ("Tauri 2.x(Rust)で Windows exe 化。GitHub Actions で自動ビルド・Release 配布", {}),
    ("コア(抽せん・リール制御)は乱数注入の純関数で、UI と完全分離", {}),
    ("開発体制", {"head": True, "bold": True, "size": 16, "color": C_GOLD}),
    ("Cursor Cloud Agent(AI エージェント)による自律開発(AGENT #001〜#098)", {}),
    ("1 サブステップ = 1 エージェント = 1 PR の分割統治", {}),
    ("引継ぎ資料(HANDOVER.md)でエージェント間の文脈を継承", {}),
], size=13))
x2 = Inches(0.35) + col_w + Inches(0.3)
add_rect(s, x2, Inches(1.25), col_w, Inches(5.7), C_PANEL)
add_text(s, x2 + Inches(0.25), Inches(1.4), col_w - Inches(0.5), Inches(5.4), bullets([
    ("AI 素材生成(fal.ai を API/MCP 経由で利用)", {"head": True, "bold": True, "size": 16, "color": C_GOLD}),
    ("画像: GPT Image 2(openai/gpt-image-2/edit)", {"bold": True}),
    ("設定資料(顔 + 全身)を参照画像に渡して演出静止画を生成", {"level": 1, "size": 12.5}),
    ("動画: Seedance 2.0(bytedance/seedance-2.0)", {"bold": True}),
    ("image-to-video / reference-to-video で演出ムービー化(顔維持ガイドライン運用)", {"level": 1, "size": 12.5}),
    ("生成ワークフロー 3 原則", {"head": True, "bold": True, "size": 16, "color": C_GOLD}),
    ("生成前にプロンプトをユーザー承認 → 生成物を確認 → 組込み", {}),
    ("指示なしの再生成禁止(API 課金の管理)", {}),
    ("原本は incoming/ へコミットして現物確保・manifest.json で出所管理", {}),
], size=13))

# ============================================================
# S4: キャラクターデザインと画像生成パイプライン
# ============================================================
s = new_slide()
header(s, "キャラクターデザインと画像生成パイプライン", "リファレンス駆動の一貫性維持")
names = ["義経", "静", "弁慶", "頼朝"]
fx = Inches(0.55)
fw = Inches(2.35)
gap = Inches(0.35)
for i, n in enumerate(names):
    p = faces.get(n)
    if p:
        add_image_fit(s, p, fx + i * (fw + gap), Inches(1.2), fw, fw)
        caption(s, fx + i * (fw + gap), Inches(3.6), fw, n + "(設定資料・顔)")
add_text(s, Inches(11.05), Inches(1.35), Inches(2.05), Inches(2.4), bullets([
    ("設定資料は 顔 + 全身 のペアで管理", {"plain": True, "size": 12}),
    ("キャラ別の禁止事項もプロンプトへ明文化", {"plain": True, "size": 12}),
], size=12))
add_rect(s, Inches(0.35), Inches(4.15), SLIDE_W - Inches(0.7), Inches(2.85), C_PANEL)
add_text(s, Inches(0.6), Inches(4.3), SLIDE_W - Inches(1.2), Inches(2.6), bullets([
    ("生成パイプライン", {"head": True, "bold": True, "size": 16, "color": C_GOLD}),
    ("① 設定資料ペア + 背景参考画像を参照に GPT Image 2 で起点静止画を生成", {}),
    ("② 連続するカットは前の生成画像を参照に連鎖 → 構図・キャラの一貫性を維持", {}),
    ("③ 静止画は 16:9 クロップ + WebP 化し「紙芝居」演出として組込み(レバーオン → 第 1 停止 → 第 3 停止で切替)", {}),
    ("④ 動画が必要な場合のみ Seedance 2.0 で image-to-video(顔ドリフト対策ガイドラインを適用)", {}),
    ("⑤ 台詞・技名・数値などの文字は画像へ焼き込まず、アプリ側でテキスト描画(差し替え耐性)", {}),
], size=13.5))

# ============================================================
# S5: 作業工程 STEP1〜6
# ============================================================
s = new_slide()
header(s, "大きな作業工程(STEP 1〜6)", "1 サブステップ = 1 PR で分割統治")
rows = [
    ("STEP 1", "リール制御", "図柄 8 種・20 コマ配列・有効 5 ライン・全役の停止制御(全 20³ 押下位置 × 押し順 6 通りの網羅テスト)", "完了(1a〜1f)", C_GREEN),
    ("STEP 2", "ゲームフロー状態遷移", "モード・前兆(偽/本)・連続演出・AT・上位 AT・エンディングの純関数ステートマシン + 100 万 G 検算", "完了(2a〜2f)", C_GREEN),
    ("STEP 3", "遊技 UI", "リール回転アニメ・タイミング目押し・押し順ナビ・メーター・SE/BGM 連動", "完了(3a〜3e)", C_GREEN),
    ("STEP 4", "演出の作り込み", "シナリオ抽せん・予告・連続演出・バトル 8G・エンディング演出(4f = 実素材差し替えが進行中)", "4a〜4e 完了", C_AMBER),
    ("STEP 5", "Windows exe 化", "Tauri 2.x 導入 + GitHub Actions で NSIS/MSI/ポータブル exe を自動ビルド", "完了(5a〜5b)", C_GREEN),
    ("STEP 6", "仕上げ", "6a 遊技データ・スランプグラフ・オートプレイ = 完了 / 6b 永続化・6c ペナルティ・6d ブランク図柄役割 = 残", "6a 完了", C_AMBER),
]
ty = Inches(1.2)
row_h = Inches(0.93)
for i, (step, name, desc, status, color) in enumerate(rows):
    y = ty + i * (row_h + Inches(0.045))
    add_rect(s, Inches(0.35), y, Inches(1.25), row_h, C_RED if i % 2 == 0 else C_PANEL2)
    add_text(s, Inches(0.35), y, Inches(1.25), row_h, [(step, 15, C_WHITE, True, 0)],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, Inches(1.65), y, Inches(9.4), row_h, C_PANEL)
    add_text(s, Inches(1.8), y + Inches(0.04), Inches(9.15), Inches(0.35),
             [(name, 13.5, C_GOLD, True, 0)])
    add_text(s, Inches(1.8), y + Inches(0.36), Inches(9.15), Inches(0.55),
             [(desc, 11, C_WHITE, False, 0)])
    add_rect(s, Inches(11.1), y, Inches(1.85), row_h, C_PANEL2)
    add_text(s, Inches(11.1), y, Inches(1.85), row_h, [(status, 12, color, True, 0)],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ============================================================
# S6: 演出素材の制作状況
# ============================================================
s = new_slide()
header(s, "演出素材の制作状況", "AI 生成画像の組込み")
pics = [
    (img_yokoku, "固有予告(静背景・強パターン)"),
    (img_battle, "上位 AT バトル(共闘)"),
    (img_battle2, "復活カットイン(静)"),
]
px = Inches(0.35)
pw = Inches(4.1)
for i, (p, cap) in enumerate(pics):
    if p:
        add_image_fit(s, p, px + i * (pw + Inches(0.17)), Inches(1.2), pw, Inches(2.45))
        caption(s, px + i * (pw + Inches(0.17)), Inches(3.62), pw, cap)
add_rect(s, Inches(0.35), Inches(4.1), SLIDE_W - Inches(0.7), Inches(2.95), C_PANEL)
add_text(s, Inches(0.6), Inches(4.22), Inches(6.1), Inches(2.7), bullets([
    ("組込み済み(AI 生成 → 紙芝居方式)", {"head": True, "bold": True, "size": 15, "color": C_GOLD}),
    ("固有予告 1(義経・静・弁慶・夕方の 4 背景 × 弱/強)", {"size": 12.5}),
    ("会話予告(固有 3。義経・静・弁慶・頼朝の全身 + 台詞)", {"size": 12.5}),
    ("下位/上位 AT バトル紙芝居(頼朝戦・8G 構成 約 50 枚)", {"size": 12.5}),
    ("エンディング・リザルト画面", {"size": 12.5}),
], size=12.5))
add_text(s, Inches(6.9), Inches(4.22), Inches(6.05), Inches(2.7), bullets([
    ("ユーザー入稿の実素材", {"head": True, "bold": True, "size": 15, "color": C_GOLD}),
    ("背景動画 7 本・図柄 8 種・AT 確定ムービー・BGM 4 曲・SE 7 音", {"size": 12.5}),
    ("残りの制作対象", {"head": True, "bold": True, "size": 15, "color": C_GOLD}),
    ("各背景の固有予告 2(2 択系)・共通予告 1・2", {"size": 12.5}),
    ("連続演出用 BGM・告知系 SE などの実素材差し替え(STEP 4f)", {"size": 12.5}),
], size=12.5))

# ============================================================
# S7: 課題
# ============================================================
s = new_slide()
header(s, "課題", "開発を通じて見えたもの")
add_rect(s, Inches(0.35), Inches(1.25), Inches(6.15), Inches(5.7), C_PANEL)
add_text(s, Inches(0.6), Inches(1.4), Inches(5.7), Inches(5.4), bullets([
    ("AI 素材生成まわり", {"head": True, "bold": True, "size": 16, "color": C_GOLD}),
    ("顔ドリフト: 生成を重ねると設定資料から顔が乖離 → 参照連鎖 + 禁止事項の明文化 + 検品で対策", {}),
    ("コンテンツモデレーション: バトル(負傷・苦痛)描写が確率的に弾かれる → 表現を弱めてリトライ", {}),
    ("API 課金: 再生成の重複コスト → プロンプト事前承認・指示なし再生成禁止をルール化", {}),
    ("動画生成は静止画より不安定 → 静止画紙芝居方式へ方針転換(品質・コスト両面で有効)", {}),
], size=13))
x2 = Inches(6.8)
add_rect(s, x2, Inches(1.25), Inches(6.15), Inches(5.7), C_PANEL)
add_text(s, x2 + Inches(0.25), Inches(1.4), Inches(5.7), Inches(5.4), bullets([
    ("開発・検証まわり", {"head": True, "bold": True, "size": 16, "color": C_GOLD}),
    ("開発 VM に GPU・音声デバイスが無く、動画再生・実音は Windows 実機確認をユーザーへ依頼", {}),
    ("エージェント間の引継ぎ: HANDOVER.md の運用で解決したが、仕様の「正」の維持に規律が必要", {}),
    ("残機能", {"head": True, "bold": True, "size": 16, "color": C_GOLD}),
    ("遊技データの永続化(6b)/ 変則押しペナルティ(6c)/ ブランク図柄の役割(6d)", {}),
    ("残り予告素材の制作・組込み(1 予告ずつ 構図承認 → 生成 → 確認 → 組込み)", {}),
], size=13))

# ============================================================
# S8: 今後の作業(応用)1 — 作業工程の短縮化
# ============================================================
s = new_slide()
header(s, "今後の作業(応用)① 作業工程の短縮化", "アニメ版権を想定した自動化")
add_rect(s, Inches(0.35), Inches(1.2), SLIDE_W - Inches(0.7), Inches(1.5), C_PANEL)
add_text(s, Inches(0.6), Inches(1.32), SLIDE_W - Inches(1.2), Inches(1.3), bullets([
    ("ゴール: アニメ化された版権機の開発で、アニメ素材(設定資料・映像・音声)を渡すだけで演出や仕様自体を自動生成する", {"plain": True, "bold": True, "size": 15, "color": C_GOLD}),
    ("今回のオリジナル版権で確立した「設定資料 → 参照連鎖 → 紙芝居組込み」パイプラインを版権素材ベースへ発展させる", {"plain": True, "size": 13}),
], size=13))
add_rect(s, Inches(0.35), Inches(2.85), Inches(6.15), Inches(3.5), C_PANEL)
add_text(s, Inches(0.6), Inches(2.97), Inches(5.7), Inches(3.3), bullets([
    ("現行の版権機開発フロー(調査)", {"head": True, "bold": True, "size": 15, "color": C_GOLD}),
    ("企画 → 絵コンテ → 素材制作 → コンポジット(AfterEffects)→ オーサリング(実機組込み)", {"size": 12}),
    ("版元からの素材支給 or 版元指定のアニメ制作会社へ新規作画を発注(スケジュール・品質のばらつきが課題)", {"size": 12}),
    ("版権キャラを崩さないことが最重要(セルルック 3DCG 化などの技法)", {"size": 12}),
    ("全カットが版元監修の往復対象 → 期間長期化の主要因", {"size": 12}),
], size=12))
x2 = Inches(6.8)
add_rect(s, x2, Inches(2.85), Inches(6.15), Inches(3.5), C_PANEL)
add_text(s, x2 + Inches(0.25), Inches(2.97), Inches(5.7), Inches(3.3), bullets([
    ("AI による短縮ポイント(構想)", {"head": True, "bold": True, "size": 15, "color": C_GOLD}),
    ("支給アニメ素材を参照画像・参照動画として GPT Image 2 / Seedance へ連鎖 → 作風を維持した新規カットを自動生成", {"size": 12}),
    ("原作の名場面・キャラ相関から、予告体系・振分けテーブルなど演出仕様自体をエージェントが自動起案", {"size": 12}),
    ("設定資料との一致検品(顔・衣装・色)を自動化し、監修往復を削減", {"size": 12}),
    ("仕様書(Excel)→ SPEC 転記 → 実装 → 網羅テストの今回の流れをテンプレ化し、機種横断で再利用", {"size": 12}),
    ("※ 版権素材の権利処理と最終監修は人間の承認が必須", {"size": 11.5, "color": C_MUTE, "plain": True}),
], size=12))
add_rect(s, Inches(0.35), Inches(6.5), SLIDE_W - Inches(0.7), Inches(0.85), C_PANEL2, line=C_GOLD)
add_text(s, Inches(0.6), Inches(6.56), SLIDE_W - Inches(1.2), Inches(0.75), bullets([
    ("日程感(目安): 生成パイプライン試作(PoC)約 1〜2 ヶ月 → 1 機種分の演出・仕様自動起案の実証 約 2〜3 ヶ月 → テンプレ化・横展開 約 1 ヶ月(計 約 4〜6 ヶ月)", {"plain": True, "bold": True, "size": 12.5, "color": C_AMBER}),
    ("※ 版元の素材支給・監修往復の待ち時間は含まない(ここが実期間の最大変動要因)", {"plain": True, "size": 10.5, "color": C_MUTE}),
], size=12))

# ============================================================
# S9: 今後の作業(応用)2 — 実機環境での動作
# ============================================================
s = new_slide()
header(s, "今後の作業(応用)② 実機環境での動作", "実機 ROM への焼き込みを目指す")
add_rect(s, Inches(0.35), Inches(1.2), Inches(6.15), Inches(5.15), C_PANEL)
add_text(s, Inches(0.6), Inches(1.32), Inches(5.7), Inches(4.9), bullets([
    ("実機の制約(規則・調査結果)", {"head": True, "bold": True, "size": 15, "color": C_GOLD}),
    ("メイン基板(出玉・リール制御): ROM 16KB(制御領域 4.5KB + データ領域 3KB)・RAM 1KB(使用 512B)", {"size": 12}),
    ("容量制約からメインプログラムは実質アセンブラで開発される", {"size": 12}),
    ("5.5 号機以降は AT 抽せん等の出玉関連処理も全てメイン基板側", {"size": 12}),
    ("サブ基板(演出・液晶): 容量制限が緩く C/C++ 等で開発。サブメイン(演出抽せん)+ サブサブ(液晶描画)の 2 層構成", {"size": 12}),
    ("型式試験: 指定試験機関(保通協 / GLI Japan)の試験 → 公安委員会の検定が必要", {"size": 12}),
], size=12))
x2 = Inches(6.8)
add_rect(s, x2, Inches(1.2), Inches(6.15), Inches(5.15), C_PANEL)
add_text(s, x2 + Inches(0.25), Inches(1.32), Inches(5.7), Inches(4.9), bullets([
    ("変換アプローチ(構想)", {"head": True, "bold": True, "size": 15, "color": C_GOLD}),
    ("① コアロジック(役抽せん・リール制御・出玉)を TypeScript → C 言語へ移植。既存の網羅テスト・100 万 G シミュレーションを同値性検証ベクタとして流用", {"size": 12}),
    ("② メイン基板相当部は容量削減の上でアセンブラ化を検討(抽せんテーブルの分母付き整数設計はそのまま活かせる)", {"size": 12}),
    ("③ 自動変換ツール(ts2c 等)は ES3 の約 70% 対応・16bit 整数のみで実用性が限定的 → AI エージェントによる移植 + テスト同値性検証が現実的", {"size": 12}),
    ("④ 演出層はデータ駆動設計(direction.ts の対応表)のため、サブ基板向け演出データテーブルへ変換しやすい", {"size": 12}),
    ("⑤ 液晶素材の実機フォーマット変換 → 評価ボードでの動作確認へ段階的に進める", {"size": 12}),
], size=12))
add_rect(s, Inches(0.35), Inches(6.45), SLIDE_W - Inches(0.7), Inches(0.9), C_PANEL2, line=C_GOLD)
add_text(s, Inches(0.6), Inches(6.51), SLIDE_W - Inches(1.2), Inches(0.8), bullets([
    ("日程感(目安): ① C 言語移植 + 同値性検証 約 2〜3 ヶ月 → ② アセンブラ化・容量最適化(16KB/1KB 適合)約 3〜6 ヶ月 → ⑤ 評価ボード動作確認まで 通算 約 6 ヶ月〜1 年", {"plain": True, "bold": True, "size": 12.5, "color": C_AMBER}),
    ("※ 実機ハード(基板・リールユニット)の調達と型式試験対応はメーカー協業が前提のため別枠(試験申請〜結果通知だけで数ヶ月規模)", {"plain": True, "size": 10.5, "color": C_MUTE}),
], size=12))

# ============================================================
# S10: まとめ
# ============================================================
s = new_slide()
header(s, "まとめ", "ロードマップ")
if img_battle:
    add_image_fit(s, img_battle, Inches(7.3), Inches(1.35), Inches(5.65), Inches(3.2))
    caption(s, Inches(7.3), Inches(4.55), Inches(5.65), "義経 × 頼朝 — 上位 AT バトル(AI 生成)")
add_text(s, Inches(0.5), Inches(1.35), Inches(6.6), Inches(5.6), bullets([
    ("現在地", {"head": True, "bold": True, "size": 17, "color": C_GOLD}),
    ("ゲーム本体(リール制御〜演出〜exe 化)は完成。演出素材の制作・組込みフェーズ", {"size": 13.5}),
    ("短期", {"head": True, "bold": True, "size": 17, "color": C_GOLD}),
    ("残り予告素材の制作・組込み / STEP 6 の永続化ほか / Windows 実機確認", {"size": 13.5}),
    ("中期(応用 ①)", {"head": True, "bold": True, "size": 17, "color": C_GOLD}),
    ("アニメ版権素材を起点にした演出・仕様の自動生成パイプライン化", {"size": 13.5}),
    ("長期(応用 ②)", {"head": True, "bold": True, "size": 17, "color": C_GOLD}),
    ("C 言語 / アセンブラへの移植と実機 ROM 化(型式試験対応を見据えた設計)", {"size": 13.5}),
], size=13.5))
add_rect(s, Inches(0.35), Inches(6.55), SLIDE_W - Inches(0.7), Pt(2.5), C_GOLD)
add_text(s, 0, Inches(6.75), SLIDE_W, Inches(0.5),
         [("AI エージェント駆動開発で「仕様 → 実装 → 素材 → 実機」までを一気通貫に", 15, C_WHITE, True, 0)],
         align=PP_ALIGN.CENTER)

out = Path(sys.argv[1]) if len(sys.argv) > 1 else ROOT / "docs/presentation/スロットアプリ開発.pptx"
out.parent.mkdir(parents=True, exist_ok=True)
prs.save(out)
print(f"saved: {out}")
